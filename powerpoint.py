
DEFAULT_NAME = "/Users/luchang/Desktop/Utilizing ChatGPT to code.pptx"

from pptx import Presentation

class PowerPoint:
    """PowerPoint class"""
    def __init__(self, filename=DEFAULT_NAME):
        """Initialize presentation with filename.
        Default /Users/luchang/Desktop/Utilizing ChatGPT to code.pptx for author use"""
        self.filename = filename
        self.presentation = Presentation(filename)
        self.title = ''
    def set_title(self, title):
        """Set title of presentation"""
        self.title = title
    def add_slide(self, text):
        """Add slide of given text"""
        layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(layout)
        shapes = slide.shapes
        shapes.placeholders[0].text_frame.text = self.title
        shapes.placeholders[1].text_frame.text = text
    def save(self, filename=None):
        """Save presentation"""
        if filename is None: filename = self.filename
        self.presentation.save(filename)

class AnimatedPowerPoint(PowerPoint):
    """Extended PowerPoint class, supports animation"""
    def __init__(self, *args, **kwargs):
        """Initialize"""
        self.CURSOR = "_"
        self.FRAMERATE = 1
        self.string = []
        self.cursor = 0
        PowerPoint.__init__(self, *args, **kwargs)
    def set_cursor(self, cursor="_"):
        """Set the cursor character to specified character"""
        self.CURSOR = cursor
    def set_framerate(self, framerate=1):
        """Set the framerate to specified framerate
        for instance, a framerate of 3 means every frame
        is 0.03 seconds."""
        self.FRAMERATE = framerate
    def write_frame(self):
        """Write the given frame with string"""
        for _ in range(self.FRAMERATE):
            self.add_slide(
                "".join(self.string[:self.cursor] +
                        [self.CURSOR] +
                        self.string[self.cursor:])
            )
    def get_string(self):
        """get current  string"""
        return self.string
    def get_cursor_index(self):
        """get current cursor index"""
        return self.cursor
    def appear(self, string):
        """Set string to string, set cursor to length of string
        then write frame"""
        self.string = list(string)
        self.cursor = len(string)
        self.write_frame()
    def disappear(self):
        """alias for appear("")"""
        self.appear("")
    def type(self, string):
        """Animates typing string at current cursor position.
        Alternatively, string can be a list of tokens"""
        for _ in string:
            self.string.insert(self.cursor, _)
            self.cursor += 1
            self.write_frame()
    def delete(self, tokens):
        """Animates deleting tokens tokens.
        IndexError if tokens exceeds limit"""
        for _ in range(tokens):
            self.cursor -= 1
            self.string.pop(self.cursor)
            self.write_frame()
    def delete_all(self):
        """alias for delete(self.cursor)"""
        self.delete(self.cursor)
    def move_left(self, tokens):
        """Animates move cursor left tokens tokens."""
        for _ in range(tokens):
            assert self.cursor != 0, "Out of range"
            self.cursor -= 1
            self.write_frame()
    def move_right(self, tokens):
        """Animates move cursor right tokens tokens."""
        for _ in range(tokens):
            assert self.cursor != len(self.string), "Out of range"
            self.cursor += 1
            self.write_frame()
    def wait(self, frames):
        """Waits for frames frames"""
        for _ in range(frames): self.write_frame()
    def blink_cursor(self, frames, cycle=None):
        """Blinks cursor through cycle for frames frames
        Then resets cursor (does not write frame)"""
        if cycle is None: cycle = [" ", self.CURSOR]
        original_cursor = self.CURSOR
        index = 0
        for _ in range(frames):
            self.set_cursor(cycle[_%len(cycle)])
            self.write_frame()
        self.set_cursor(original_cursor)
